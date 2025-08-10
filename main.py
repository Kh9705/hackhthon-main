# main.py
import httpx
import fitz  # PyMuPDF
import numpy as np
from typing import List, Dict, Any
from fastapi import FastAPI, HTTPException
from pydantic import BaseModel, HttpUrl
from sentence_transformers import SentenceTransformer
import faiss
from dotenv import load_dotenv
import os
import time
import torch
import asyncio
import io
import re
from pptx.enum.shapes import MSO_SHAPE_TYPE

from pptx import Presentation
import pandas as pd
from PIL import Image
import pytesseract
import zipfile
import magic

from urllib.parse import urlparse

# --- MODIFIED: Import only the Mistral client ---
from mistralai.async_client import MistralAsyncClient

# Advanced text splitter
from langchain.text_splitter import RecursiveCharacterTextSplitter

# --- NEW: Import libraries for DOCX and Email processing ---
import docx
import email
from email.policy import default
import extract_msg

# --- Environment and API Setup ---
load_dotenv()

# --- MODIFIED: Load multiple Mistral API keys ---
API_KEYS = [
    key for key in os.environ if key.startswith("MISTRAL_API_KEY_")
]
if not API_KEYS:
    print("[WARNING] No MISTRAL_API_KEY_<n> environment variables set!")
    MISTRAL_CREDENTIALS = []
else:
    MISTRAL_CREDENTIALS = [os.getenv(key) for key in API_KEYS]
    print(f"[INFO] Loaded {len(MISTRAL_CREDENTIALS)} Mistral API keys.")

# Create a list of async Mistral clients, one for each key
ASYNC_CLIENTS = [
    MistralAsyncClient(api_key=key) for key in MISTRAL_CREDENTIALS
]

# --- NEW: Dynamic Chunking & Retrieval Configuration ---
# Thresholds for document size based on character count
SMALL_DOC_THRESHOLD = 20000  # characters
LARGE_DOC_THRESHOLD = 100000 # characters

# Configuration map: size -> (chunk_size, chunk_overlap, k_for_retrieval)
CHUNK_CONFIG = {
    "small":  (500, 75, 10),
    "medium": (1000, 150, 7),
    "large":  (2000, 300, 5)
}

# --- FastAPI App Initialization ---
app = FastAPI(
    title="Optimized Multi-Format RAG Service with Mistral AI",
    description="RAG service for PDF, DOCX, and Email files using Mistral AI.",
)

# --- Pydantic Models (Unchanged) ---
class QuestionPayload(BaseModel):
    documents: HttpUrl
    questions: List[str]

class AnswerOut(BaseModel):
    answers: List[str]

# --- In-Memory Cache (MODIFIED for clarity) ---
document_cache: Dict[str, Dict[str, Any]] = {}

# --- Device Selection (Unchanged) ---
device = "cuda" if torch.cuda.is_available() else "cpu"
print(f"[INFO] Using device for embeddings: {device}")

# --- Load SentenceTransformer Model (Unchanged) ---
print("[INFO] Loading SentenceTransformer model: all-MiniLM-L6-v2...")
encoder = SentenceTransformer("all-MiniLM-L6-v2", device=device)
print("[INFO] Model loaded successfully.")


# --- NEW: Helper function to extract text from EML files ---
# --- MODIFIED: A robust new function to handle raw email content ---
def _extract_text_from_eml(eml_bytes: bytes) -> str:
    """
    Decodes the entire raw email content, including all headers,
    into a single string. This allows for analysis of headers, metadata, and body.
    It tries common encodings to prevent errors.
    """
    try:
        # Try decoding with UTF-8 first, the most common encoding.
        return eml_bytes.decode('utf-8')
    except UnicodeDecodeError:
        # If UTF-8 fails, fall back to latin-1, which is also common and less prone to errors.
        print("[WARNING] Could not decode email as UTF-8, falling back to latin-1.")
        return eml_bytes.decode('latin-1', errors='ignore')


# --- MODIFIED: Core Logic generalized for multiple document types and dynamic chunking ---
# --- MODIFIED: Core Logic with a robust fix for URL parsing ---
# --- MODIFIED: Core Logic with a robust fix for URL parsing ---

def process_and_index_document(doc_url: str):
    """
    Downloads, extracts text, and indexes any supported document type.
    (PDF, DOCX, EML, MSG, PPTX, XLSX, PNG, JPEG, ZIP)
    """
    global document_cache
    print(f"[INFO] Processing document: {doc_url}")

    # Use regex to robustly find the file extension
    match = re.search(r"\.(pdf|docx|eml|msg|pptx|xlsx|png|jpeg|jpg|zip|bin)(?=\?|$)", doc_url, re.IGNORECASE)
    file_extension = match.group(0).lower() if match else None
    
    try:
        # 1. Download Document
        with httpx.Client() as http_client:
            response = http_client.get(doc_url, timeout=120.0) # Increased timeout for large files
            response.raise_for_status()
        doc_bytes = response.content
        print(f"[INFO] Document downloaded ({len(doc_bytes)} bytes).")

        # 2. Extract text based on file type
        full_text = ""
        print(f"[INFO] Detected file extension: '{file_extension}'")
        
        if file_extension in [".pdf", ".docx", ".eml", ".msg"]:
            # Re-use your existing logic for these types
            if file_extension == ".pdf":
                doc = fitz.open("pdf", doc_bytes)
                full_text = "".join(page.get_text() for page in doc)
                doc.close()
            elif file_extension == ".docx":
                doc_stream = io.BytesIO(doc_bytes)
                document = docx.Document(doc_stream)
                full_text = "\n".join([para.text for para in document.paragraphs])
            elif file_extension == ".eml":
                full_text = _extract_text_from_eml(doc_bytes) # Your existing helper
            elif file_extension == ".msg":
                msg = extract_msg.Message(doc_bytes)
                full_text = msg.body
        
        # --- NEW: Handle PPTX files ---
        elif file_extension == ".pptx":
            print("[INFO] Processing PPTX file, checking for text and images...")
            ppt_stream = io.BytesIO(doc_bytes)
            prs = Presentation(ppt_stream)
            
            # Loop through each slide to find text and images
            for slide_number, slide in enumerate(prs.slides):
                print(f"[INFO]  - Processing Slide {slide_number + 1}...")
                for shape in slide.shapes:
                    # 1. Extract text from normal text boxes
                    if shape.has_text_frame:
                        full_text += shape.text_frame.text + "\n"
                    
                    # 2. Extract text from images using OCR
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        print(f"[INFO]    - Found an image, running OCR...")
                        try:
                            image = shape.image
                            image_bytes = image.blob
                            img = Image.open(io.BytesIO(image_bytes))
                            ocr_text = pytesseract.image_to_string(img)
                            if ocr_text.strip():
                                full_text += f"--- OCR Text from Image ---\n{ocr_text}\n--- End of OCR Text ---\n"
                        except Exception as ocr_error:
                            print(f"[WARNING] Could not process an image on slide {slide_number + 1}: {ocr_error}")


        # --- NEW: Handle XLSX files ---
        elif file_extension == ".xlsx":
            xls_stream = io.BytesIO(doc_bytes)
            # Read all sheets into a dictionary of DataFrames
            sheets = pd.read_excel(xls_stream, sheet_name=None)
            for sheet_name, df in sheets.items():
                full_text += f"--- Sheet: {sheet_name} ---\n"
                full_text += df.to_string() + "\n\n"
        
        # --- NEW: Handle Image files with OCR ---
        elif file_extension in [".png", ".jpeg", ".jpg"]:
            image_stream = io.BytesIO(doc_bytes)
            image = Image.open(image_stream)
            full_text = pytesseract.image_to_string(image)
            print("[INFO] Extracted text from image using OCR.")
            
        # --- NEW: Handle ZIP archives ---
        elif file_extension == ".zip":
            zip_stream = io.BytesIO(doc_bytes)
            archive = zipfile.ZipFile(zip_stream)
            for filename in archive.namelist():
                if not filename.endswith('/'): # Ignore directories
                    print(f"[INFO] Processing '{filename}' from ZIP archive...")
                    file_bytes = archive.read(filename)
                    # Use python-magic to detect the inner file type
                    inner_file_type = magic.from_buffer(file_bytes, mime=True)
                    # This is a simplified handler; a full implementation might recursively call this function
                    if 'pdf' in inner_file_type:
                        doc = fitz.open("pdf", file_bytes)
                        full_text += f"--- Content of {filename} ---\n" + "".join(page.get_text() for page in doc) + "\n\n"
                        doc.close()
                    # Add more handlers for other file types inside zips if needed
                    else:
                        try: # Try to decode as text as a fallback
                           full_text += f"--- Content of {filename} ---\n" + file_bytes.decode('utf-8', errors='ignore') + "\n\n"
                        except Exception:
                           print(f"[WARNING] Could not extract text from '{filename}' inside ZIP.")
        
        # --- NEW: Handle unsupported binary files ---
        elif file_extension == ".bin":
            raise HTTPException(status_code=415, detail="Unsupported file type: .bin files cannot be processed for text.")
            
        else:
            raise HTTPException(status_code=415, detail=f"Unsupported file type: '{file_extension}'.")
            
        # --- The rest of your function remains the same ---
        print(f"[INFO] Extracted {len(full_text)} characters.")
        # ... (if not full_text.strip(): ... etc)
        # ... (text_splitter = ...)
        # ... (vectors = encoder.encode(...) etc.)

    # ... (the rest of your function, including the except block)


        if not full_text.strip():
            print("[WARNING] No text extracted.")
            document_cache[doc_url] = {"index": None, "chunks": [], "timestamp": time.time()}
            return

        # 3. Determine document size and select chunking strategy
        doc_len = len(full_text)
        if doc_len < SMALL_DOC_THRESHOLD:
            size_category = "small"
        elif doc_len > LARGE_DOC_THRESHOLD:
            size_category = "large"
        else:
            size_category = "medium"

        chunk_size, chunk_overlap, k_value = CHUNK_CONFIG[size_category]
        print(f"[INFO] Document size is '{size_category}'. Using chunk_size={chunk_size}, k={k_value}")

        # 4. Split text into chunks with the chosen strategy
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            length_function=len,
        )
        chunks = text_splitter.split_text(full_text)
        print(f"[INFO] Created {len(chunks)} chunks.")

        # 5. Create embeddings on GPU (Unchanged)
        print("[INFO] Generating embeddings on GPU...")
        vectors = encoder.encode(
            chunks,
            batch_size=128,
            show_progress_bar=True,
            convert_to_numpy=True
        ).astype("float32")

        # 6. Create FAISS index on CPU (Unchanged)
        print("[INFO] Creating FAISS index on CPU...")
        index = faiss.IndexFlatL2(vectors.shape[1])
        index.add(vectors)
        print("[INFO] FAISS CPU index created.")

        # 7. Store in cache
        document_cache[doc_url] = {
            "index": index,
            "chunks": chunks,
            "timestamp": time.time(),
            "k_value": k_value
        }
        print(f"[SUCCESS] Indexed: {doc_url}")

    except Exception as e:
        print(f"[ERROR] Document processing failed: {e}")
        if isinstance(e, HTTPException):
            raise
        raise HTTPException(status_code=500, detail=f"Failed to process document: {e}")


# --- MODIFIED: Helper function for a single async Mistral LLM call ---
async def get_answer_from_llm(client: MistralAsyncClient, messages: List[Dict[str, str]], question: str) -> str:
    """Makes a single API call to the Mistral LLM and handles errors."""
    try:
        response = await client.chat(
            model="mistral-large-latest",
            messages=messages,
            temperature=0.1,
            max_tokens=500,
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        print(f"[ERROR] Mistral LLM call for question '{question}' failed: {e}")
        return "Error while generating answer."


# --- API Endpoint (MODIFIED for dynamic k retrieval) ---
@app.post("/hackrx/run", response_model=AnswerOut)
async def run_rag(payload: QuestionPayload):
    doc_url_str = str(payload.documents)
    questions = payload.questions
    print(f"[INFO] Request for {doc_url_str} with {len(questions)} questions.")

    if not ASYNC_CLIENTS:
        raise HTTPException(status_code=500, detail="No Mistral API keys configured.")

    if doc_url_str not in document_cache:
        print("[INFO] Document not in cache. Processing...")
        process_and_index_document(doc_url_str)
    else:
        print("[INFO] Using cached document index.")

    cached_data = document_cache[doc_url_str]
    index = cached_data["index"]
    chunks = cached_data["chunks"]

    if not chunks or index is None:
        return AnswerOut(
            answers=["I could not extract any text from the provided document."] * len(questions)
        )
    
    # --- MODIFIED: Use the dynamic K value from the cache ---
    # Use .get() with a default for backward compatibility with any old cache entries
    default_k = CHUNK_CONFIG["medium"][2]
    k = cached_data.get("k_value", default_k)
    print(f"[INFO] Using k={k} for retrieval based on document size.")

    tasks = []
    for i, question in enumerate(questions):
        print(f"[INFO] Preparing task for: '{question}'")

        question_vector = encoder.encode(
            [question], convert_to_numpy=True
        ).astype("float32")
        
        # Use the dynamically determined 'k' for the search
        _, indices = index.search(question_vector, k)
        context = "\n\n---\n\n".join(chunks[i] for i in indices[0])

        prompt = f"""
You are an expert Q&A system. Use the context to answer accurately and concisely.
If the answer is not in the context, say you cannot answer.

## Context:
{context}

## Question:
{question}

## Answer:
"""
        messages = [{"role": "user", "content": prompt.strip()}]
        
        client_to_use = ASYNC_CLIENTS[i % len(ASYNC_CLIENTS)]
        
        task = get_answer_from_llm(client_to_use, messages, question)
        tasks.append(task)
        
        await asyncio.sleep(0.5)

    print(f"[INFO] Running {len(tasks)} tasks in parallel with Mistral AI...")
    final_answers = await asyncio.gather(*tasks)
    print("[INFO] All tasks completed.")

    return AnswerOut(answers=final_answers)

# --- MODIFIED: Health Check to reflect new cache name ---
@app.get("/")
def health_check():
    return {"status": "ok", "cached_docs": list(document_cache.keys())}